"""
PPT 确认单生成引擎

生成标准5页确认单：
1. 封面 - 项目名称、户型面积、方案编号、商家Logo/名称(会员)
2. 铺贴效果图 - 排版渲染图
3. 材料明细 - 主砖+辅料列表
4. 商家信息 - (会员专享)
5. 确认签字区
"""
import io
import os
import tempfile
from typing import Optional, List, Dict, Any
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from datetime import datetime


BRAND_PRIMARY = RGBColor(0x1A, 0x36, 0x5D)
BRAND_ACCENT = RGBColor(0xD4, 0xA5, 0x74)
BRAND_LIGHT = RGBColor(0xE8, 0xED, 0xF2)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK = RGBColor(0x33, 0x33, 0x33)
GRAY = RGBColor(0x99, 0x99, 0x99)
LIGHT_GRAY = RGBColor(0xF0, 0xF0, 0xF0)
TABLE_HEADER_BG = RGBColor(0x1A, 0x36, 0x5D)
TABLE_HEADER_FG = RGBColor(0xFF, 0xFF, 0xFF)
TABLE_ALT_BG = RGBColor(0xF7, 0xF9, 0xFC)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def _add_brand_bar(slide, top: float = 0):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(top), SLIDE_W, Inches(0.06)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = BRAND_ACCENT
    shape.line.fill.background()


def _add_footer(slide):
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(7.0), Inches(12), Inches(0.3))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "本方案由排砖宝 TileLayout AI 生成  |  www.tilelayout.ai"
    p.font.size = Pt(8)
    p.font.color.rgb = GRAY
    p.alignment = PP_ALIGN.CENTER


def _add_logo_or_placeholder(slide, is_member: bool, store_name: str = "", logo_path: str = ""):
    if is_member and store_name:
        txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(4), Inches(0.5))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = store_name
        p.font.size = Pt(20)
        p.font.color.rgb = BRAND_PRIMARY
        p.font.bold = True
        if logo_path and os.path.exists(logo_path):
            slide.shapes.add_picture(logo_path, Inches(0.8), Inches(0.4), Inches(0.8), Inches(0.8))
    else:
        txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.7), Inches(5), Inches(0.4))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = "⬆ 升级会员，展示您的品牌"
        p.font.size = Pt(14)
        p.font.color.rgb = GRAY
        p.font.italic = True


def create_confirmation_ppt(
    project_data: Dict[str, Any],
    is_member: bool = False,
    store_profile: Optional[Dict[str, Any]] = None,
    layout_preview_image: Optional[bytes] = None,
    materials: Optional[List[Dict[str, Any]]] = None,
    auxiliary_materials: Optional[Dict[str, Any]] = None,
    show_price: bool = True,
    output_path: Optional[str] = None,
) -> io.BytesIO:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    blank_layout = prs.slide_layouts[6]

    project_name = project_data.get("project_name", "未命名方案")
    project_area = project_data.get("area_sq_m", 0)
    project_id = project_data.get("project_id", "NEW-001")
    stats = project_data.get("statistics", {})
    store = store_profile or {}

    # ==================== 第1页：封面 ====================
    slide1 = prs.slides.add_slide(blank_layout)
    bg_shape = slide1.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_W, SLIDE_H
    )
    bg_shape.fill.solid()
    bg_shape.fill.fore_color.rgb = WHITE
    bg_shape.line.fill.background()

    _add_brand_bar(slide1, top=0)
    _add_brand_bar(slide1, top=7.44)

    _add_logo_or_placeholder(slide1, is_member, store.get("store_name", ""))

    txBox = slide1.shapes.add_textbox(Inches(2), Inches(2.5), Inches(9), Inches(1.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "瓷砖铺贴方案确认单"
    p.font.size = Pt(40)
    p.font.color.rgb = BRAND_PRIMARY
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    p2 = tf.add_paragraph()
    p2.text = project_name
    p2.font.size = Pt(24)
    p2.font.color.rgb = BRAND_ACCENT
    p2.alignment = PP_ALIGN.CENTER
    p2.space_before = Pt(16)

    info_text = f"户型面积: {project_area}m²  |  方案编号: {project_id}  |  {datetime.now().strftime('%Y年%m月%d日')}"
    txBox2 = slide1.shapes.add_textbox(Inches(2), Inches(4.5), Inches(9), Inches(0.5))
    tf2 = txBox2.text_frame
    p3 = tf2.paragraphs[0]
    p3.text = info_text
    p3.font.size = Pt(14)
    p3.font.color.rgb = GRAY
    p3.alignment = PP_ALIGN.CENTER

    _add_footer(slide1)

    # ==================== 第2页：铺贴效果图 ====================
    slide2 = prs.slides.add_slide(blank_layout)
    slide2.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_W, SLIDE_H
    ).fill.solid()
    slide2.shapes[0].fill.fore_color.rgb = WHITE
    slide2.shapes[0].line.fill.background()
    _add_brand_bar(slide2, top=0)

    if layout_preview_image:
        img_stream = io.BytesIO(layout_preview_image)
        slide2.shapes.add_picture(img_stream, Inches(1), Inches(0.8), Inches(11), Inches(5.5))
    else:
        placeholder = slide2.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(2), Inches(2), Inches(9), Inches(3)
        )
        placeholder.fill.solid()
        placeholder.fill.fore_color.rgb = LIGHT_GRAY
        placeholder.line.color.rgb = GRAY
        tf = placeholder.text_frame
        tf.paragraphs[0].text = "铺贴效果图"
        tf.paragraphs[0].font.size = Pt(24)
        tf.paragraphs[0].font.color.rgb = GRAY
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    txBox = slide2.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.5))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = f"铺贴效果图 - {project_name}"
    p.font.size = Pt(18)
    p.font.color.rgb = BRAND_PRIMARY
    p.font.bold = True

    _add_footer(slide2)

    # ==================== 第3页：材料明细 ====================
    slide3 = prs.slides.add_slide(blank_layout)
    slide3.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_W, SLIDE_H
    ).fill.solid()
    slide3.shapes[0].fill.fore_color.rgb = WHITE
    slide3.shapes[0].line.fill.background()
    _add_brand_bar(slide3, top=0)

    txBox = slide3.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.5))
    p = txBox.text_frame.paragraphs[0]
    p.text = f"材料明细清单 - {project_name}"
    p.font.size = Pt(20)
    p.font.color.rgb = BRAND_PRIMARY
    p.font.bold = True

    all_materials = []
    if materials:
        all_materials.extend(materials)
    if auxiliary_materials and auxiliary_materials.get("cost_items"):
        all_materials.extend(auxiliary_materials["cost_items"])

    if all_materials:
        cols = 6 if (show_price and is_member) else 4
        rows = len(all_materials) + 1

        col_widths = [Inches(1.5), Inches(2.5), Inches(2), Inches(2), Inches(1.5), Inches(1.8)] if cols == 6 else [Inches(1.8), Inches(3.5), Inches(3), Inches(3.2)]
        headers = ["品名", "规格", "数量", "单位", "单价(元)", "金额(元)"] if cols == 6 else ["品名", "规格", "数量", "单位"]

        table_left = Inches(0.5)
        table_top = Inches(1.2)
        table_w = sum(col_widths)
        table_shape = slide3.shapes.add_table(rows, cols, table_left, table_top, table_w, Inches(0.45 * rows))
        table = table_shape.table

        for ci, (cw, header) in enumerate(zip(col_widths, headers)):
            table.columns[ci].width = cw
            cell = table.cell(0, ci)
            cell.text = header
            cell.fill.solid()
            cell.fill.fore_color.rgb = TABLE_HEADER_BG
            for pp in cell.text_frame.paragraphs:
                pp.font.size = Pt(11)
                pp.font.color.rgb = TABLE_HEADER_FG
                pp.font.bold = True
                pp.alignment = PP_ALIGN.CENTER

        for ri, mat in enumerate(all_materials):
            for ci, key in enumerate(["name", "spec", "qty", "unit", "unit_price", "amount"] if cols == 6 else ["name", "spec", "qty", "unit"]):
                cell = table.cell(ri + 1, ci)
                val = mat.get(key, "")
                if isinstance(val, (int, float)):
                    cell.text = str(val) if key == "qty" else f"¥{val:.2f}"
                else:
                    cell.text = str(val) if val is not None else "-"
                if ri % 2 == 1:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = TABLE_ALT_BG
                for pp in cell.text_frame.paragraphs:
                    pp.font.size = Pt(10)
                    pp.alignment = PP_ALIGN.CENTER if ci >= 2 else PP_ALIGN.LEFT

        if show_price and is_member:
            total_val = sum(m.get("amount", 0) for m in all_materials)
            total_top = table_top + Inches(0.45 * (rows + 1))
            txBox = slide3.shapes.add_textbox(Inches(8), total_top, Inches(4.5), Inches(0.4))
            p = txBox.text_frame.paragraphs[0]
            p.text = f"合计: ¥{total_val:.2f}"
            p.font.size = Pt(16)
            p.font.color.rgb = BRAND_PRIMARY
            p.font.bold = True
            p.alignment = PP_ALIGN.RIGHT

    _add_footer(slide3)

    # ==================== 第4页：商家信息 ====================
    slide4 = prs.slides.add_slide(blank_layout)
    slide4.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_W, SLIDE_H
    ).fill.solid()
    slide4.shapes[0].fill.fore_color.rgb = WHITE
    slide4.shapes[0].line.fill.background()
    _add_brand_bar(slide4, top=0)

    txBox = slide4.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.5))
    p = txBox.text_frame.paragraphs[0]
    p.text = "商家信息与售后保障"
    p.font.size = Pt(20)
    p.font.color.rgb = BRAND_PRIMARY
    p.font.bold = True

    if is_member and store.get("store_name"):
        info_items = [
            ("门店名称", store.get("store_name", "")),
            ("联系电话", store.get("phone", "")),
            ("门店地址", store.get("address", "")),
        ]
        for i, (label, val) in enumerate(info_items):
            txBox = slide4.shapes.add_textbox(Inches(2), Inches(2 + i * 1), Inches(9), Inches(0.6))
            tf = txBox.text_frame
            p = tf.paragraphs[0]
            p.text = f"{label}: {val}"
            p.font.size = Pt(18)
            p.font.color.rgb = DARK
    else:
        box = slide4.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(2), Inches(2), Inches(9), Inches(3)
        )
        box.fill.solid()
        box.fill.fore_color.rgb = LIGHT_GRAY
        box.line.color.rgb = GRAY
        tf = box.text_frame
        tf.word_wrap = True
        p1 = tf.paragraphs[0]
        p1.text = "🔓 升级为付费会员，展示您的品牌信息"
        p1.font.size = Pt(24)
        p1.font.color.rgb = BRAND_ACCENT
        p1.alignment = PP_ALIGN.CENTER
        p2 = tf.add_paragraph()
        p2.text = "包含门店Logo、名称、联系方式、地址"
        p2.font.size = Pt(14)
        p2.font.color.rgb = GRAY
        p2.alignment = PP_ALIGN.CENTER
        p2.space_before = Pt(16)

    _add_footer(slide4)

    # ==================== 第5页：确认签字 ====================
    slide5 = prs.slides.add_slide(blank_layout)
    slide5.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_W, SLIDE_H
    ).fill.solid()
    slide5.shapes[0].fill.fore_color.rgb = WHITE
    slide5.shapes[0].line.fill.background()
    _add_brand_bar(slide5, top=0)

    txBox = slide5.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.5))
    p = txBox.text_frame.paragraphs[0]
    p.text = "客户确认签字"
    p.font.size = Pt(20)
    p.font.color.rgb = BRAND_PRIMARY
    p.font.bold = True

    info_text = (
        "本人已确认上述铺贴方案、材料清单及费用明细（如有），同意按此方案进行施工。\n\n"
        "备注说明：________________________________________\n\n"
        "客户签字：________________     日期：____年____月____日\n\n"
        "设计师签字：________________     日期：____年____月____日"
    )
    txBox2 = slide5.shapes.add_textbox(Inches(2), Inches(2.5), Inches(9), Inches(3.5))
    tf = txBox2.text_frame
    tf.word_wrap = True
    p2 = tf.paragraphs[0]
    p2.text = info_text
    p2.font.size = Pt(16)
    p2.font.color.rgb = DARK
    p2.line_spacing = Pt(28)

    _add_footer(slide5)

    if output_path:
        prs.save(output_path)
        return io.BytesIO()

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf
